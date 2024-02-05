import zipfile
from flask import Flask, request, jsonify
import csv
from pptx import Presentation
import PyPDF2
import olefile
import re
import tempfile
import io
from docx import Document
from PIL import Image
import pyclamd

app = Flask(__name__)

@app.route('/scan', methods=['POST'])
def scan():
    if request.method == 'POST':
        file = request.files['file']
        filename = file.filename
        file_extension = filename.split('.')[-1].lower()

        if file_extension in ['zip', 'pdf', 'pptx', 'csv', 'xslx', 'docx']:
            try:
                if file_extension == 'zip':
                    is_malicious, message = zipscan(file)
                elif file_extension == 'pdf':
                    is_malicious, message = pdfscan(file)
                elif file_extension == 'pptx':
                    is_malicious, message = pptxscan(file)
                elif file_extension == 'docx':
                    is_malicious, message = docxscan(file)
                elif file_extension == 'xslx':
                    is_malicious, message = xlsxscan(file)
                elif file_extension == 'csv':
                    is_malicious, message = csvscan(file)
                elif file_extension == 'png':
                    is_malicious, message = imagescan(file)
                else:
                    is_malicious = False
                    message = "File type not currently supported for malicious content analysis."
            except Exception as e:
                is_malicious = False
                message = f"Error analyzing the file: {str(e)}"

            return jsonify({'is_malicious': is_malicious, 'message': message})
        else:
            return jsonify({'error': 'Unsupported file type'})

    return jsonify({'error': 'Invalid request'})


# ScanMethods
def imagescan(file):
    try:
        cd = pyclamd.ClamdUnixSocket()
        scan_result = cd.scan_file(file)

        if scan_result[file] == 'OK':
            return False, "The image appears to be valid and does not contain known indicators of malicious content."
        else:
            return True, f"The image may contain malicious content: {scan_result[file]}"

    except Exception as e:
        return True, f"Error analyzing the image: {str(e)}"

def zipscan(file):
    try:
        with zipfile.ZipFile(file, 'r') as zip_file:
            for file_info in zip_file.infolist():
                file_name = file_info.filename
                file_content = zip_file.read(file_info)

                # Check for common malicious patterns in file names and content
                if file_name.lower().endswith('.exe'):
                    return True, f"The ZIP file likely contains an executable malware program: {file_name}"
                elif file_name.lower().endswith(('.js')):
                    return True, f"The ZIP file contains a potentially dangerous script: {file_name}"
                elif b'<script>' in file_content or b'<object>' in file_content:  # Check for embedded scripts
                    return True, f"The ZIP file contains a file with potential script code: {file_name}"
                elif file_name.lower().endswith(('.vbs')):
                    return True, f"The ZIP file contains a potentially dangerous VBscript: {file_name}"
                elif file_name.lower().endswith(('.bat')):
                    return True, f"The ZIP file contains a potentially dangerous Bat script: {file_name}"
                elif file_name.lower().endswith(('.cmd')):
                    return True, f"The ZIP file contains a potentially dangerous CMD script: {file_name}"

    except Exception as e:
        return False, f"Error analyzing the ZIP file: {str(e)}"
    return False, "The ZIP file does not seem to contain known indicators of malicious content."


def pdfscan(file):
    try:
        with file.stream as pdf_file:
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            # Check each page for common malicious patterns
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                page_text = page.extract_text()

                # Check for common malicious patterns
                if "malicious_pattern_1" in page_text.lower():
                    return True, "The PDF contains a known malicious pattern: malicious_pattern_1"
                elif "malicious_pattern_2" in page_text.lower():
                    return True, "The PDF contains a known malicious pattern: malicious_pattern_2"
                # Add more checks as needed

            return False, "The PDF does not seem to contain known indicators of malicious content."

    except Exception as e:
        return False, f"Error analyzing the PDF: {str(e)}"



def pptxscan(file):
    try:
        prs = Presentation(io.BytesIO(file.read()))

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    text = shape.text_frame.text
                    if "javascript:" in text.lower() or "vbscript:" in text.lower():
                        return True, "The PPTX contains potential script code in a text frame."
    except Exception as e:
        return False, f"Error analyzing the PPTX file: {str(e)}"
    return False, "The PPTX file does not seem to contain known indicators of malicious content."


def docxscan(file):
    try:
        doc = Document(file)
        for rel in doc.part.rels.values():
            if "macro" in rel.reltype:
                return True, "The DOCX file contains macros, which could be malicious."

    except Exception as e:
        return False, f"Error analyzing the DOCX file: {str(e)}"
    
    return False, "The DOCX file does not seem to contain known indicators of malicious content."

def xlsxscan(file):
    try:
        ole = olefile.OleFileIO(io.BytesIO(file.read()))

        for stream in ole.listdir():
            if stream[6:] == 'Macros':  # Check for macros
                return True, "The XLSX file contains macros, which could be malicious."

    except Exception as e:
        return False, f"Error analyzing the XLSX file: {str(e)}"
    return False, "The XLSX file does not seem to contain known indicators of malicious content."


def csvscan(file):
    try:
        # Get the stream from the FileStorage object
        csv_file = io.StringIO(file.stream.read().decode('utf-8'))
        csv_reader = csv.reader(csv_file)

        for row in csv_reader:
            for cell in row:
                # Check for common malicious patterns in cell content
                if any(pattern in cell.lower() for pattern in ['javascript:', 'vbscript:', 'eval(', 'cmd.exe', 'powershell']):
                    return True, f"The CSV file contains potential script code or suspicious commands in a cell: {cell}"
                if len(cell) > 1000 and re.match(r'^[a-zA-Z0-9+/]+={0,2}$', cell):  # Check for long encoded strings
                    return True, f"The CSV file contains a potentially encoded string: {cell}"
                if re.match(r'\\\\[^\\]+\\[^\\]+\\.exe$', cell):  # Check for unusual file paths
                    return True, f"The CSV file contains a suspicious file path: {cell}"

        return False, "The CSV file does not seem to contain known indicators of malicious content."
    except Exception as e:
        return False, f"Error analyzing the CSV file: {str(e)}"


if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
